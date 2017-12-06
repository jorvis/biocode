#!/usr/bin/env python3

"""

Custom request from our Director on High to provide a script which generates
a PowerPoint slide set based on a specific input convention in a text file.

INPUT:

- Spaces at the beginning of a line (0-2) are header levels
- Content is composed entirely of bulleted lists
- List elements start with # symbols, and more than one indicates nesting levels

Example:

This would be title 1
#this would be text for first bullet
#this would be text for second bullet
##text for sub-bullet
##second sub-bullet
#text for third bullet

This would be title 2
#blah1a
#blah1b
##blah2

Module documentation:
https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

"""

import argparse
import os
from pptx import Presentation

def main():
    parser = argparse.ArgumentParser( description='Create PowerPoint slides from a plain text file')

    parser.add_argument('-i', '--input_file', type=str, required=True, help='Path to an input text file to be read' )
    parser.add_argument('-o', '--output_file', type=str, required=True, help='Path to an output pptx file to be created' )
    args = parser.parse_args()

    prs = Presentation()
    bullet_layout = prs.slide_layouts[1]

    current_slide = None
    current_body_shape = None
    current_tf = None

    for line in open(args.input_file):
        line = line.rstrip()

        # skip blank lines
        if len(line) == 0:
            continue

        if line.startswith('   '):
            raise Exception("Header levels greater than 3 are not supported")
        elif line.startswith('  '):
            raise Exception('Sub-headers not yet supported')
        elif line.startswith(' '):
            raise Exception('Sub-headers not yet supported')
        elif line.startswith('#'):
            hash_count = get_hash_count(line)

            if hash_count == 1:
                if current_tf is None:
                    current_tf = current_body_shape.text_frame
                    current_tf.text = line.lstrip('#')
                else:
                    p = current_tf.add_paragraph()
                    p.text = line.lstrip('#')
                    p.level = hash_count - 1
            elif hash_count > 1:
                p = current_tf.add_paragraph()
                p.text = line.lstrip('#')
                p.level = hash_count - 1
            
        else:
            # Must be a title slide line
            current_slide = prs.slides.add_slide(bullet_layout)
            shapes = current_slide.shapes

            title_shape = shapes.title
            title_shape.text = line
            
            current_body_shape = shapes.placeholders[1]
            current_tf = None

    prs.save(args.output_file)


def get_hash_count(s):
    cnt = 0
    for c in s:
        if c == '#':
            cnt += 1
        else:
            break

    return cnt

if __name__ == '__main__':
    main()







